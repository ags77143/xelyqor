import os
import re
import tempfile
from typing import Optional


def extract_youtube_transcript(url: str) -> str:
    """Try youtube-transcript-api first, fall back to yt-dlp + Whisper."""
    video_id = _extract_video_id(url)
    if not video_id:
        raise ValueError("Could not extract YouTube video ID from URL.")

    # Attempt 1: youtube-transcript-api
    try:
        from youtube_transcript_api import YouTubeTranscriptApi, NoTranscriptFound, TranscriptsDisabled
        
        transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
        
        # Priority order for captions
        transcript = None
        try:
            transcript = transcript_list.find_manually_created_transcript(['en'])
        except Exception:
            pass
        
        if not transcript:
            try:
                transcript = transcript_list.find_generated_transcript(['en'])
            except Exception:
                pass
        
        if not transcript:
            # Try any manual transcript
            for t in transcript_list:
                if not t.is_generated:
                    transcript = t
                    break
        
        if not transcript:
            # Try any transcript at all
            for t in transcript_list:
                transcript = t
                break
        
        if transcript:
            fetched = transcript.fetch()
            return " ".join(entry["text"] for entry in fetched)
            
    except Exception as e:
        print(f"youtube-transcript-api failed: {e}")

    # Attempt 2: yt-dlp + Whisper
    return _whisper_transcribe(url)


def _extract_video_id(url: str) -> Optional[str]:
    patterns = [
        r'(?:v=|/v/|youtu\.be/|/embed/)([a-zA-Z0-9_-]{11})',
    ]
    for p in patterns:
        m = re.search(p, url)
        if m:
            return m.group(1)
    return None


def _whisper_transcribe(url: str) -> str:
    try:
        import yt_dlp
        import whisper

        with tempfile.TemporaryDirectory() as tmpdir:
            audio_path = os.path.join(tmpdir, "audio.mp3")
            
            ydl_opts = {
                'format': 'bestaudio/best',
                'outtmpl': os.path.join(tmpdir, 'audio.%(ext)s'),
                'postprocessors': [{
                    'key': 'FFmpegExtractAudio',
                    'preferredcodec': 'mp3',
                    'preferredquality': '128',
                }],
                'quiet': True,
            }
            
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                ydl.download([url])
            
            # Find the audio file
            for f in os.listdir(tmpdir):
                if f.endswith('.mp3'):
                    audio_path = os.path.join(tmpdir, f)
                    break
            
            model = whisper.load_model("base")
            result = model.transcribe(audio_path)
            return result["text"]
            
    except Exception as e:
        raise ValueError(
            f"Could not extract transcript from YouTube video. "
            f"Error: {str(e)}\n\n"
            f"You can manually get a transcript from:\n"
            f"• https://tactiq.io — Paste the YouTube URL to get a transcript\n"
            f"• https://downsub.com — Download subtitles as text"
        )


def extract_pdf_text(file_bytes: bytes) -> str:
    import fitz  # PyMuPDF
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
    return text


def extract_pptx_text(file_bytes: bytes) -> str:
    import io
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text)
    return "\n\n".join(text_parts)
